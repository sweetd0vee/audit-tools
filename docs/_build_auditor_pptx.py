"""Build the Audit Copilot presentation. Run: python docs/_build_auditor_pptx.py"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import nsmap, qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

# --- palette ---
NAVY = RGBColor(0x0B, 0x1F, 0x33)
NAVY2 = RGBColor(0x12, 0x2C, 0x45)
INK = RGBColor(0x14, 0x22, 0x33)
CREAM = RGBColor(0xF6, 0xF1, 0xE6)
PAPER = RGBColor(0xFF, 0xFC, 0xF7)
GOLD = RGBColor(0xC4, 0x96, 0x2C)
GOLD2 = RGBColor(0xE8, 0xC9, 0x6A)
TEAL = RGBColor(0x2F, 0x8F, 0x84)
TEAL_DK = RGBColor(0x1E, 0x5F, 0x58)
BRICK = RGBColor(0xB4, 0x4A, 0x3A)
MUTED = RGBColor(0x6B, 0x7C, 0x8D)
SLATE = RGBColor(0x3D, 0x4F, 0x61)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0x17, 0x33, 0x4C)
LINE = RGBColor(0xE4, 0xD9, 0xC4)

W, H = Inches(13.333), Inches(7.5)
FONT = "Calibri"
FONT_B = "Calibri"


def _set_run(run, text, size, color, bold=False, italic=False, font=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    # Force latin + ea so PowerPoint on Windows doesn't swap to Times
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", font)


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _fill_line(shape, fill, line=None, width_pt=1.0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(width_pt)


def _rect(slide, l, t, w, h, color, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    _fill_line(s, color, line)
    s.shadow.inherit = False
    return s


def _round(slide, l, t, w, h, color, line=None, adj=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    _fill_line(s, color, line)
    try:
        s.adjustments[0] = adj
    except Exception:
        pass
    s.shadow.inherit = False
    return s


def _tb(shape, text, size, color, bold=False, align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.clear()
    run = p.add_run()
    _set_run(run, text, size, color, bold, italic)
    return tf


def _add_p(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT, italic=False, space_before=6, space_after=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    _set_run(run, text, size, color, bold, italic)
    return p


def _textbox(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    _set_run(run, text, size, color, bold, italic)
    return box, tf


def _footer(slide, n, total, dark=False):
    c = RGBColor(0x9A, 0xB0, 0xC0) if dark else MUTED
    _textbox(slide, Inches(0.5), Inches(7.18), Inches(8.5), Inches(0.28),
             "Аудитор  ·  коробочный ИИ-агент внутренней проверки  ·  конфиденциально",
             10, c, italic=True)
    _textbox(slide, Inches(11.4), Inches(7.18), Inches(1.4), Inches(0.28),
             f"{n}  /  {total}", 10, c, align=PP_ALIGN.RIGHT)


def _kicker(slide, text, l=Inches(0.55), t=Inches(0.28), dark=False):
    bar = _rect(slide, l, t + Inches(0.07), Inches(0.18), Inches(0.18), GOLD)
    col = GOLD2 if dark else GOLD
    _textbox(slide, l + Inches(0.28), t, Inches(8), Inches(0.32), text.upper(), 11, col, bold=True)


def _title(slide, text, l=Inches(0.55), t=Inches(0.52), w=Inches(12.2), size=28, color=INK):
    _textbox(slide, l, t, w, Inches(0.7), text, size, color, bold=True)


def new_light(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, W, H, PAPER)
    _rect(s, 0, 0, W, Inches(0.08), NAVY)
    _rect(s, 0, Inches(7.08), W, Inches(0.42), CREAM)
    _rect(s, 0, Inches(7.08), W, Pt(1.5), GOLD)
    return s


def new_dark(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, W, H, NAVY)
    _rect(s, 0, Inches(7.08), W, Inches(0.42), NAVY2)
    _rect(s, 0, Inches(7.08), W, Pt(1.5), GOLD)
    return s


def card(slide, l, t, w, h, fill=WHITE, line=LINE):
    return _round(slide, l, t, w, h, fill, line, adj=0.06)


def chevron(slide, l, t, w, h, fill, text, tsize=11, tcolor=WHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, l, t, w, h)
    _fill(s, fill)
    s.shadow.inherit = False
    _tb(s, text, tsize, tcolor, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s


TOTAL = 19


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # ========== 1 TITLE ==========
    s = new_dark(prs)
    _rect(s, 0, 0, Inches(0.18), H, GOLD)
    _textbox(s, Inches(0.7), Inches(1.35), Inches(12), Inches(0.35),
             "ВНУТРЕННИЙ АУДИТ  ·  БАНК РБ  ·  ОН-ПРЕМ  ·  OPEN WEBUI PIPE",
             12, GOLD, bold=True)
    _textbox(s, Inches(0.7), Inches(1.85), Inches(12), Inches(1.3),
             "Аудитор", 54, WHITE, bold=True)
    _textbox(s, Inches(0.7), Inches(3.15), Inches(11.5), Inches(1.1),
             "ИИ-агент сопровождения внутренней аудиторской проверки.\nНе чат с моделью — цикл с руками, стоп-краном и системой записи.",
             20, RGBColor(0xD5, 0xE0, 0xEA))
    # chips
    chips = [
        ("Коробка", "compose up → :3000"),
        ("Pipe, не ReAct", "фаза выбирается кодом"),
        ("HITL", "без «утверждаю» — нет скачивания"),
        ("On-prem", "данные не уходят в облако"),
    ]
    for i, (a, b) in enumerate(chips):
        x = Inches(0.7) + Inches(i * 3.1)
        c = _round(s, x, Inches(4.7), Inches(2.9), Inches(1.15), CARD, adj=0.1)
        tf = _tb(c, a, 14, GOLD, bold=True)
        tf.margin_left = Inches(0.18)
        tf.margin_top = Inches(0.16)
        _add_p(tf, b, 12, RGBColor(0xC5, 0xD4, 0xE0), space_before=4)
    _footer(s, 1, TOTAL, dark=True)

    # ========== 2 AGENDA ==========
    s = new_light(prs)
    _kicker(s, "Содержание")
    _title(s, "О чём эта презентация")
    items = [
        ("01", "Зачем", "Боль аудитора, инвариант продукта, почему не «ещё один чат»."),
        ("02", "История и функционал", "Сеанс проверки, команды, артефакты Word / Excel / zip."),
        ("03", "Open WebUI Pipe", "Теория Functions, реализация class Pipe, почему не native tools."),
        ("04", "Алгоритм", "Маршрутизация фаз кодом, HITL, RAG с цитатой из утверждённого акта."),
        ("05", "Коробка и контур", "On-prem, allowlist РБ, конфиденциальные данные, поставка."),
        ("06", "Эффект", "Что автоматизируется сегодня и где человек остаётся обязательным."),
    ]
    for i, (n, h, d) in enumerate(items):
        col, row = i % 3, i // 3
        x = Inches(0.5) + Inches(col * 4.2)
        y = Inches(1.45) + Inches(row * 2.55)
        c = card(s, x, y, Inches(4.0), Inches(2.35))
        _textbox(s, x + Inches(0.22), y + Inches(0.22), Inches(3.5), Inches(0.4), n, 20, GOLD, bold=True)
        _textbox(s, x + Inches(0.22), y + Inches(0.7), Inches(3.5), Inches(0.4), h, 18, INK, bold=True)
        _textbox(s, x + Inches(0.22), y + Inches(1.2), Inches(3.5), Inches(0.9), d, 13, SLATE)
    _footer(s, 2, TOTAL)

    # ========== 3 PROBLEM ==========
    s = new_light(prs)
    _kicker(s, "Контекст")
    _title(s, "Каждую проверку аудитор собирает заново")
    pains = [
        ("01", "Вспомнить НПА", "ГК, НК, инструкции НБРБ, акты Минфина и МНС — по теме этой недели."),
        ("02", "Найти редакцию", "pravo.by / nbrb.by, актуальный текст, не методичка с форума."),
        ("03", "Вычитать кодекс", "Десятки страниц, чтобы вытащить 5–7 норм в рабочую бумагу."),
        ("04", "Спроектировать проверку", "Программа, гипотезы, критерии — вручную, каждый раз с нуля."),
    ]
    for i, (n, h, d) in enumerate(pains):
        y = Inches(1.4) + Inches(i * 1.15)
        c = card(s, Inches(0.5), y, Inches(8.15), Inches(1.05))
        num = _round(s, Inches(0.7), y + Inches(0.26), Inches(0.55), Inches(0.52), NAVY, adj=0.15)
        _tb(num, n, 12, GOLD, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _textbox(s, Inches(1.45), y + Inches(0.16), Inches(7), Inches(0.35), h, 16, INK, bold=True)
        _textbox(s, Inches(1.45), y + Inches(0.52), Inches(7), Inches(0.4), d, 13, SLATE)
    q = _round(s, Inches(8.9), Inches(1.4), Inches(3.9), Inches(4.6), NAVY, adj=0.05)
    tf = _tb(q, "РАЗРЫВ", 11, GOLD, bold=True)
    tf.margin_left = Inches(0.28)
    tf.margin_top = Inches(0.3)
    _add_p(tf, "Модель без документов\nгаллюцинирует статьи.", 16, WHITE, bold=True, space_before=14)
    _add_p(tf, "Документы без модели\nне масштабируются.", 16, WHITE, bold=True, space_before=16)
    _add_p(tf, "Нужен RAG + доменный сервер: модель думает, НПА доказывают, аудитор подтверждает.", 13, RGBColor(0xC5, 0xD4, 0xE0), space_before=22)
    _footer(s, 3, TOTAL)

    # ========== 4 PRODUCT ==========
    s = new_light(prs)
    _kicker(s, "Продукт")
    _title(s, "Коробка уровня Cursor — для внутренней проверки")
    _textbox(s, Inches(0.55), Inches(1.25), Inches(12.2), Inches(0.55),
             "Cursor победил не тем, что написал IDE с нуля. Он взял VS Code и сделал законченный продукт. Мы делаем то же со стеком Open WebUI + Ollama + SearXNG + тонкий Audit Tool Server.",
             15, SLATE)
    rows = [
        ("Cursor для кода", "Аудитор для проверки"),
        ("Поставил — открыл — пишешь", "Поставил — открыл — ведёшь проверку"),
        ("VS Code не форкали", "Open WebUI не форкаем"),
        ("Цитата из этого репозитория", "Цитата из этого утверждённого акта"),
        ("Один вход: окно редактора", "Один вход: чат на :3000"),
    ]
    hdr = _rect(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.48), NAVY)
    _textbox(s, Inches(0.75), Inches(2.04), Inches(5.5), Inches(0.32), "Cursor", 14, WHITE, bold=True)
    _textbox(s, Inches(6.85), Inches(2.04), Inches(5.5), Inches(0.32), "Мы", 14, GOLD, bold=True)
    for i, (a, b) in enumerate(rows[1:], start=0):
        y = Inches(2.43) + Inches(i * 0.7)
        bg = CREAM if i % 2 == 0 else WHITE
        r = _rect(s, Inches(0.5), y, Inches(12.3), Inches(0.7), bg, LINE)
        _textbox(s, Inches(0.75), y + Inches(0.16), Inches(5.6), Inches(0.4), a, 15, INK)
        _rect(s, Inches(6.55), y + Inches(0.18), Pt(2), Inches(0.34), GOLD)
        _textbox(s, Inches(6.85), y + Inches(0.16), Inches(5.7), Inches(0.4), b, 15, INK, bold=True)
    _footer(s, 4, TOTAL)

    # ========== 5 INVARIANT ==========
    s = new_dark(prs)
    _kicker(s, "Инвариант продукта", dark=True)
    _title(s, "Модель не ставит суждение", color=WHITE, size=30)
    q = _round(s, Inches(0.55), Inches(1.45), Inches(12.2), Inches(2.15), CARD, adj=0.04)
    tf = _tb(q, "«Модель ускоряет поиск нормы и черновик обзора. Она не ставит аудиторское суждение и не заменяет чтение первоисточника. Цитата проверяется по файлу в кейсе.»",
             20, WHITE, italic=True, anchor=MSO_ANCHOR.MIDDLE)
    tf.margin_left = Inches(0.45)
    tf.margin_right = Inches(0.45)
    points = [
        ("Ускоряет", "Библиотека НПА, саммари, программа, гипотезы, вопрос по статье."),
        ("Не подписывает", "Нет вывода «нарушение / не нарушение». Нет автоподписи WP."),
        ("Сверяется", "Блок «Откуда в базе знаний» + zip первоисточника в папке кейса."),
    ]
    for i, (h, d) in enumerate(points):
        x = Inches(0.55) + Inches(i * 4.15)
        c = _round(s, x, Inches(3.9), Inches(3.95), Inches(2.0), CARD, adj=0.08)
        tf = _tb(c, h, 16, GOLD, bold=True)
        tf.margin_left = Inches(0.22)
        tf.margin_top = Inches(0.22)
        _add_p(tf, d, 13, RGBColor(0xC5, 0xD4, 0xE0), space_before=10)
    _footer(s, 5, TOTAL, dark=True)

    # ========== 6 USER STORY ==========
    s = new_light(prs)
    _kicker(s, "Пользовательская история")
    _title(s, "Внутренний аудитор банка, проверка аренды")
    _textbox(s, Inches(0.55), Inches(1.22), Inches(12.2), Inches(0.4),
             "Тема недели: аренда коммерческой недвижимости — валютные расчёты, НДС, договоры с нерезидентами.",
             14, SLATE, italic=True)

    left = card(s, Inches(0.5), Inches(1.75), Inches(6.0), Inches(4.95))
    tf = _tb(left, "РАНЬШЕ", 12, BRICK, bold=True)
    tf.margin_left = Inches(0.28)
    tf.margin_top = Inches(0.22)
    for line in [
        "Вспоминаю акты по памяти и старым WP.",
        "Ищу редакции на pravo.by вкладками.",
        "Читаю ГК/НК целиком «на всякий случай».",
        "Выписываю нормы в черновик вручную.",
        "Программу проверки собираю с нуля.",
        "Гипотезы — в голове или в Excel с нуля.",
        "Риск: устаревшая редакция или чужая юрисдикция.",
    ]:
        _add_p(tf, "▸  " + line, 14, INK, space_before=10)

    right = _round(s, Inches(6.8), Inches(1.75), Inches(6.0), Inches(4.95), NAVY, adj=0.05)
    tf = _tb(right, "ТЕПЕРЬ", 12, GOLD, bold=True)
    tf.margin_left = Inches(0.28)
    tf.margin_top = Inches(0.22)
    for line in [
        "Одной фразой описываю проверку в чате.",
        "Модель предлагает 8–15 актов — я утверждаю номера.",
        "Официальные тексты качаются в папку кейса.",
        "саммари — Word на 6–10 стр. со ссылками.",
        "программа проверки — черновик СВА в Word.",
        "гипотезы — Excel-чеклист 8–10 пунктов.",
        "вопрос … — норма с цитатой из этой библиотеки.",
    ]:
        _add_p(tf, "▸  " + line, 14, WHITE, space_before=10)
    _footer(s, 6, TOTAL)

    # ========== 7 JOURNEY ==========
    s = new_light(prs)
    _kicker(s, "Сценарий сеанса")
    _title(s, "Один чат = одна проверка")
    steps = [
        ("1", "Описать", "Проверка аренды…,\nаренда, валюта, НДС"),
        ("2", "Утвердить", "утверждаю 1, 2, 4\nплюс Инструкция 38"),
        ("3", "Скачать", "pravo.by / nbrb.by\nиндекс по статьям"),
        ("4", "Читать", "саммари\nпрограмма проверки"),
        ("5", "Планировать", "гипотезы\nExcel-чеклист"),
        ("6", "Спросить", "вопрос Какой срок\nрегистрации аренды?"),
    ]
    for i, (n, h, d) in enumerate(steps):
        x = Inches(0.4) + Inches(i * 2.15)
        if i:
            arr = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, x - Inches(0.22), Inches(2.35), Inches(0.2), Inches(0.16)
            )
            _fill(arr, GOLD)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), Inches(1.45), Inches(0.55), Inches(0.55))
        _fill(circ, NAVY)
        _tb(circ, n, 16, GOLD, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        c = card(s, x, Inches(2.15), Inches(2.02), Inches(2.55))
        tf = _tb(c, h, 14, INK, bold=True, align=PP_ALIGN.CENTER)
        tf.margin_top = Inches(0.22)
        for part in d.split("\n"):
            _add_p(tf, part, 11, SLATE, align=PP_ALIGN.CENTER, space_before=6)

    note = card(s, Inches(0.5), Inches(4.95), Inches(12.3), Inches(1.75))
    tf = _tb(note, "Правила сеанса", 14, GOLD, bold=True)
    tf.margin_left = Inches(0.28)
    tf.margin_top = Inches(0.18)
    _add_p(tf, "Модель в шапке — «Аудитор» (Pipe), не qwen3.8:27b с иконкой ol. Новая тема — новый чат. Без слова «вопрос» ответ идёт из знаний модели, не из библиотеки. Пока нет «утверждаю» — ничего не скачивается.",
           14, INK, space_before=8)
    _footer(s, 7, TOTAL)

    # ========== 8 FEATURES ==========
    s = new_light(prs)
    _kicker(s, "Функционал")
    _title(s, "Что умеет агент сегодня")
    feats = [
        ("Библиотека НПА", "Propose → HITL → download с allowlist РБ. extra_titles, ручной URL, known_sources."),
        ("Саммари Word", "Карточка по каждому акту + обзор проверки + приложение фрагментов [n]."),
        ("Саммари total", "Конспект «из головы» модели — не цитата. Явно отделён от базы."),
        ("Программа проверки", "Черновик СВА: цель, риски, процедуры, критерии по праву РБ."),
        ("Гипотезы Excel", "8–10 гипотез для проверки. Опирается на саммари / программу / НПА."),
        ("Вопрос по базе", "Префикс «вопрос …» → RAG + блок «Откуда в базе знаний»."),
        ("Обычный диалог", "Без префикса — POST /chat. Без цитат. Не тащить в WP как норму."),
        ("Система записи", "Кейс на диске: статусы, манифест, sha256, zip первоисточников."),
    ]
    for i, (h, d) in enumerate(feats):
        col, row = i % 4, i // 4
        x = Inches(0.4) + Inches(col * 3.22)
        y = Inches(1.4) + Inches(row * 2.6)
        c = card(s, x, y, Inches(3.08), Inches(2.4))
        bar = _rect(s, x, y, Inches(0.12), Inches(2.4), GOLD if row == 0 else TEAL)
        _textbox(s, x + Inches(0.28), y + Inches(0.22), Inches(2.65), Inches(0.7), h, 15, INK, bold=True)
        _textbox(s, x + Inches(0.28), y + Inches(0.95), Inches(2.65), Inches(1.25), d, 12, SLATE)
    _footer(s, 8, TOTAL)

    # ========== 9 ARTIFACTS ==========
    s = new_light(prs)
    _kicker(s, "Выходы")
    _title(s, "Артефакты, которые аудитор уносит из чата")
    arts = [
        (".zip", "Библиотека НПА", "Первоисточники HTML/PDF, которые вы утвердили — не «то, что модель вспомнила»."),
        (".docx", "Обзор актов", "Нормативный контур проверки, карточки существенного, ссылки на статьи."),
        (".docx", "Программа", "Черновик программы аудиторской проверки банка РБ. Не подпись руководителя СВА."),
        (".xlsx", "Гипотезы", "Чеклист 8–10 гипотез: что смотреть, на что опираться, статус для HITL."),
        (".docx", "Total", "Конспект из знаний модели. Полезно как ориентир, не как норма."),
        ("чат", "Цитата", "Ответ на «вопрос …» + блок источников. Сверяется с файлом в zip."),
    ]
    for i, (ext, h, d) in enumerate(arts):
        col, row = i % 3, i // 3
        x = Inches(0.45) + Inches(col * 4.25)
        y = Inches(1.4) + Inches(row * 2.55)
        c = card(s, x, y, Inches(4.05), Inches(2.35))
        badge = _round(s, x + Inches(0.22), y + Inches(0.22), Inches(0.95), Inches(0.38), NAVY, adj=0.3)
        _tb(badge, ext, 11, GOLD, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _textbox(s, x + Inches(1.3), y + Inches(0.24), Inches(2.5), Inches(0.38), h, 16, INK, bold=True)
        _textbox(s, x + Inches(0.22), y + Inches(0.8), Inches(3.6), Inches(1.3), d, 13, SLATE)
    _footer(s, 9, TOTAL)

    # ========== 10 ARCHITECTURE ==========
    s = new_light(prs)
    _kicker(s, "Архитектура")
    _title(s, "Сборка готового, не форк и не свой чат")
    layers = [
        (NAVY, WHITE, GOLD, "Поверхность", "Open WebUI — чат, auth, файлы, статусы. Свой React не пишем."),
        (TEAL_DK, WHITE, GOLD2, "Клей агента", "Function Pipe «Аудитор». Фазы кодом. Без бизнес-логики внутри."),
        (RGBColor(0x3A, 0x2A, 0x18), WHITE, GOLD, "Ядро", "Audit Tool Server (FastAPI): кейс, HITL, добыча НПА, ingest, RAG, Word/Excel."),
        (SLATE, WHITE, GOLD2, "Инфра", "Ollama (LLM + embed) · SearXNG (allowlist) · диск кейса. DuckDB — план."),
    ]
    for i, (bg, tc, ac, h, d) in enumerate(layers):
        y = Inches(1.38) + Inches(i * 1.25)
        c = _round(s, Inches(0.5), y, Inches(12.3), Inches(1.12), bg, adj=0.04)
        tf = _tb(c, h.upper(), 12, ac, bold=True)
        tf.margin_left = Inches(0.35)
        tf.margin_top = Inches(0.16)
        _add_p(tf, d, 16, tc, space_before=6)
    _footer(s, 10, TOTAL)

    # ========== 11 PIPE THEORY ==========
    s = new_light(prs)
    _kicker(s, "Теория  ·  Open WebUI")
    _title(s, "Три способа «сделать агента» — выбран Pipe")
    cols = [
        ("Chat model", "Обычный Qwen в списке (иконка ol).",
         "Свободный диалог. Нет кейса, нет стоп-крана, нет download.",
         "Не продукт", BRICK),
        ("Tools + Native FC", "Модель сама вызывает HTTP-tools.",
         "Локальная 27B плохо держит function calling: может скачать без «утверждаю» или не вызвать tool вовсе.",
         "Запасной путь", MUTED),
        ("Function Pipe", "Python class Pipe в Admin → Functions.",
         "Сообщение входит в pipe(). Фазу выбирает код. Модель вызывается точечно. HITL железобетонный.",
         "Продуктовый путь", TEAL),
    ]
    for i, (h, sub, body, tag, tagc) in enumerate(cols):
        x = Inches(0.4) + Inches(i * 4.25)
        c = card(s, x, Inches(1.4), Inches(4.1), Inches(5.3))
        _rect(s, x, Inches(1.4), Inches(4.1), Inches(0.12), tagc)
        _textbox(s, x + Inches(0.25), Inches(1.7), Inches(3.6), Inches(0.45), h, 18, INK, bold=True)
        _textbox(s, x + Inches(0.25), Inches(2.2), Inches(3.6), Inches(0.7), sub, 13, SLATE, italic=True)
        _textbox(s, x + Inches(0.25), Inches(2.95), Inches(3.6), Inches(2.4), body, 14, INK)
        badge = _round(s, x + Inches(0.25), Inches(5.85), Inches(2.4), Inches(0.45), tagc, adj=0.3)
        _tb(badge, tag, 12, WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _footer(s, 11, TOTAL)

    # ========== 12 PIPE IMPLEMENTATION ==========
    s = new_light(prs)
    _kicker(s, "Реализация  ·  audit_agent.py v0.2.7")
    _title(s, "Как устроен Function Pipe «Аудитор»")

    left = card(s, Inches(0.45), Inches(1.38), Inches(6.3), Inches(5.3))
    tf = _tb(left, "Контракт Open WebUI", 15, GOLD, bold=True)
    tf.margin_left = Inches(0.28)
    tf.margin_top = Inches(0.2)
    for line in [
        "class Pipe — появляется в списке как модель «Аудитор».",
        "Valves — параметры без правки кода: AUDIT_API, PUBLIC_API, таймауты, ключ Knowledge.",
        "async pipe(body, __user__, __event_emitter__) — единственная точка входа.",
        "__event_emitter__ рисует статусы в чате («Скачиваю…»), это не ответ модели.",
        "Бизнес-логики в Pipe нет: только маршрут + HTTP на FastAPI + формат ответа.",
        "Память проверки: скрытая метка <!--audit-case:id--> в истории чата.",
    ]:
        _add_p(tf, "●  " + line, 13, INK, space_before=11)

    right = _round(s, Inches(6.95), Inches(1.38), Inches(5.9), Inches(5.3), NAVY, adj=0.04)
    tf = _tb(right, "Valves (из коробки)", 15, GOLD, bold=True)
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.22)
    for k, v in [
        ("AUDIT_API", "http://backend:8100"),
        ("PUBLIC_API", "http://localhost:8100"),
        ("TIMEOUT_SEC", "600  — propose / download"),
        ("BRIEF_TIMEOUT_SEC", "1800 — Word / Excel"),
        ("OPENWEBUI_API_KEY", "пусто = индекс сервера"),
    ]:
        _add_p(tf, k, 13, GOLD2, bold=True, space_before=14)
        _add_p(tf, v, 13, WHITE, space_before=2)
    _footer(s, 12, TOTAL)

    # ========== 13 ALGORITHM ==========
    s = new_light(prs)
    _kicker(s, "Алгоритм")
    _title(s, "Фаза выбирается кодом. Первое совпадение побеждает")
    _textbox(s, Inches(0.55), Inches(1.22), Inches(12.2), Inches(0.35),
             "Это не свободный ReAct. Локальная 27B не решает, какой tool вызвать.",
             14, SLATE, italic=True)
    steps = [
        ("1", "помощь / ?", "шпаргалка, без LLM"),
        ("2", "вопрос … / /ask", "RAG по индексу кейса"),
        ("3", "утверждаю гипотезы", "HITL (задел под мнение)"),
        ("4", "мнение / заключение", "Word раздела I и полный черновик АЗ"),
        ("5", "программа проверки", "черновик СВА → .docx"),
        ("6", "саммари total", "конспект из знаний модели"),
        ("7", "гипотезы", "Excel-чеклист 8–10"),
        ("8", "саммари", "карточки актов → .docx"),
        ("9", "утверждаю / скачай", "select + download, без LLM"),
        ("10", "документы / статус", "чтение кейса, без LLM"),
        ("11", "похож на проверку", "create + propose, если нет кейса"),
        ("12", "иначе", "обычный чат POST /chat"),
    ]
    for i, (n, h, d) in enumerate(steps):
        col, row = i % 4, i // 4
        x = Inches(0.4) + Inches(col * 3.22)
        y = Inches(1.65) + Inches(row * 1.7)
        c = card(s, x, y, Inches(3.08), Inches(1.52))
        num = _round(s, x + Inches(0.16), y + Inches(0.18), Inches(0.42), Inches(0.38), NAVY, adj=0.2)
        _tb(num, n, 11, GOLD, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _textbox(s, x + Inches(0.68), y + Inches(0.18), Inches(2.2), Inches(0.4), h, 13, INK, bold=True)
        _textbox(s, x + Inches(0.16), y + Inches(0.7), Inches(2.75), Inches(0.65), d, 12, SLATE)
    _footer(s, 13, TOTAL)

    # ========== 14 HITL ==========
    s = new_light(prs)
    _kicker(s, "Human-in-the-loop")
    _title(s, "Стоп-кран: модель предлагает, аудитор решает")
    flow = [
        ("LLM", "JSON со списком\n8–15 актов"),
        ("Pipe", "Нумерует, просит\n«утверждаю»"),
        ("Человек", "Номера, плюс,\nURL, отказ"),
        ("Сервер", "select → download\nтолько после HITL"),
        ("Индекс", "Статьи, embeddings,\nбиблиотека кейса"),
    ]
    for i, (h, d) in enumerate(flow):
        x = Inches(0.4) + Inches(i * 2.55)
        if i:
            a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x - Inches(0.28), Inches(2.35), Inches(0.24), Inches(0.18))
            _fill(a, GOLD)
        bg = NAVY if i in (0, 3, 4) else (TEAL if i == 1 else GOLD)
        tc = INK if i == 2 else WHITE
        c = _round(s, x, Inches(1.5), Inches(2.35), Inches(2.05), bg, adj=0.08)
        tf = _tb(c, h, 16, GOLD2 if i != 2 else NAVY, bold=True, align=PP_ALIGN.CENTER)
        tf.margin_top = Inches(0.28)
        for part in d.split("\n"):
            _add_p(tf, part, 13, tc, align=PP_ALIGN.CENTER, space_before=6)

    notes = [
        ("Где модель молчит", "утверждаю, плюс, скачай, документы, статус, помощь — чистый код. Галлюцинация статьи на этих шагах невозможна."),
        ("Где модель говорит", "propose (JSON), саммари, total, программа, гипотезы, вопрос по фрагментам, обычный чат."),
        ("Политика download", "URL: ручная ссылка → known_sources → SearXNG. Домены только РБ. Клиентский текст в поиск не идёт."),
    ]
    for i, (h, d) in enumerate(notes):
        x = Inches(0.4) + Inches(i * 4.25)
        c = card(s, x, Inches(3.85), Inches(4.1), Inches(2.8))
        tf = _tb(c, h, 14, GOLD, bold=True)
        tf.margin_left = Inches(0.22)
        tf.margin_top = Inches(0.2)
        _add_p(tf, d, 13, INK, space_before=10)
    _footer(s, 14, TOTAL)

    # ========== 15 RAG ==========
    s = new_light(prs)
    _kicker(s, "Качество ответа")
    _title(s, "Цитата из утверждённого акта, не из памяти модели")
    left = card(s, Inches(0.45), Inches(1.38), Inches(6.35), Inches(5.3))
    tf = _tb(left, "Конвейер «вопрос …»", 15, GOLD, bold=True)
    tf.margin_left = Inches(0.28)
    tf.margin_top = Inches(0.2)
    for n, line in [
        ("1", "Только если сообщение начинается с «вопрос» — иначе это не норма."),
        ("2", "Эмбеддинг вопроса той же моделью, что индекс: qwen3-embedding."),
        ("3", "Hybrid: BM25 + пересечение слов + cosine → RRF → MMR. Top-k = 16."),
        ("4", "Промпт: только фрагменты. Нет во фрагментах — отказ. Не РФ / IFRS."),
        ("5", "Pipe дописывает блок «Откуда в базе знаний» даже если модель гладкая."),
        ("6", "Веб-поиска на этом шаге нет. SearXNG — только добыча в библиотеку."),
    ]:
        _add_p(tf, f"{n}.  {line}", 13, INK, space_before=12)

    right_top = _round(s, Inches(7.0), Inches(1.38), Inches(5.85), Inches(2.45), NAVY, adj=0.05)
    tf = _tb(right_top, "Норма  ≠  факт", 16, GOLD, bold=True)
    tf.margin_left = Inches(0.28)
    tf.margin_top = Inches(0.22)
    _add_p(tf, "Норма — НПА в knowledge_*. Ответ: RAG + цитата статьи.", 13, WHITE, space_before=10)
    _add_p(tf, "Факт — выгрузки клиента. Ответ: SQL / tools (план). Не класть Excel в Knowledge.", 13, WHITE, space_before=8)

    right_bot = card(s, Inches(7.0), Inches(4.05), Inches(5.85), Inches(2.63))
    tf = _tb(right_bot, "Нарезка под юридический текст", 14, GOLD, bold=True)
    tf.margin_left = Inches(0.22)
    tf.margin_top = Inches(0.18)
    _add_p(tf, "Чанкер по «Статья N», не дешёвый сплиттер по символам. Окно модели 32k. Embedding — не MiniLM. Температура 0.1–0.2.", 13, INK, space_before=8)
    _add_p(tf, "Источник истины для аудитора — блок цитат и файл в zip, не красноречие абзаца.", 13, INK, space_before=8)
    _footer(s, 15, TOTAL)

    # ========== 16 CONFIDENTIAL ==========
    s = new_light(prs)
    _kicker(s, "Конфиденциальность")
    _title(s, "Работа с чувствительными данными — by design")
    items = [
        ("On-prem LLM", "Ollama на хосте банка. qwen3.8:27b не уходит в OpenAI / облако. Формулировки проверки остаются в контуре."),
        ("Allowlist РБ", "pravo.by, nbrb.by, minfin, nalog, government, president. Defense in depth: SearXNG и сервер."),
        ("Клиент ≠ поиск", "Выгрузки, договоры, суммы в SearXNG не передаются. Веб-поиск в чате выключен."),
        ("Две вселенные", "Норма и факт не смешиваются в одной Knowledge. Скрепка с Excel клиента — запрещённый путь."),
        ("Air-gap готовность", "Образы и веса моделей едут файлами, не качаются с Hugging Face в банке."),
        ("Аудит действий", "case.json, manifest (URL, sha256), статусы created → ready. Trail вызовов — в плане."),
    ]
    for i, (h, d) in enumerate(items):
        col, row = i % 3, i // 3
        x = Inches(0.4) + Inches(col * 4.25)
        y = Inches(1.4) + Inches(row * 2.55)
        c = card(s, x, y, Inches(4.1), Inches(2.35))
        ico = _round(s, x + Inches(0.22), y + Inches(0.22), Inches(0.16), Inches(0.16), TEAL, adj=0.5)
        _textbox(s, x + Inches(0.5), y + Inches(0.16), Inches(3.35), Inches(0.45), h, 16, INK, bold=True)
        _textbox(s, x + Inches(0.22), y + Inches(0.75), Inches(3.65), Inches(1.4), d, 13, SLATE)
    _footer(s, 16, TOTAL)

    # ========== 17 BOX ==========
    s = new_light(prs)
    _kicker(s, "Поставка")
    _title(s, "Коробочное решение, не стенд разработчика")
    _textbox(s, Inches(0.55), Inches(1.22), Inches(12.2), Inches(0.4),
             "Если аудитор должен выбрать embedding, включить hybrid и привязать Knowledge — это не продукт.",
             14, SLATE, italic=True)

    cmd = _round(s, Inches(0.5), Inches(1.75), Inches(12.3), Inches(1.15), NAVY, adj=0.05)
    tf = _tb(cmd, "docker compose up -d --build     →     http://localhost:3000     →     модель «Аудитор»",
             18, WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    box_items = [
        ("Уже в коробке", [
            "Open WebUI + backend + SearXNG в одном compose",
            "Ollama на хосте: чат 27B + embedding",
            "RAG-env из compose: hybrid, qwen embedding, top-k 16, num_ctx 32k",
            "Веб-поиск выключен, имя продукта — «Аудитор»",
            "Pipe — один Python-файл, без форка Open WebUI",
        ]),
        ("Мнение из коробки", [
            "Промпты методолога ВА банка РБ",
            "Запрет РФ / ЕС / IFRS, если их нет во фрагментах",
            "Температура 0.1–0.2 на сервере, не слайдер чата",
            "Шаблон RAG: нет во фрагментах — отказ",
            "HITL на список НПА железный",
        ]),
        ("Что это даёт банку", [
            "Один URL для аудитора, без Vite и Swagger",
            "Свой код только там, где рынка нет: кейс, домен РБ",
            "Можно поставить на машину / контур СВА",
            "Не SaaS с клиентскими файлами во внешней LLM",
            "Масштаб: следующий аудитор открывает тот же чат",
        ]),
    ]
    for i, (h, lines) in enumerate(box_items):
        x = Inches(0.45) + Inches(i * 4.25)
        c = card(s, x, Inches(3.15), Inches(4.1), Inches(3.55))
        tf = _tb(c, h, 15, GOLD, bold=True)
        tf.margin_left = Inches(0.22)
        tf.margin_top = Inches(0.18)
        for line in lines:
            _add_p(tf, "·  " + line, 12, INK, space_before=7)
    _footer(s, 17, TOTAL)

    # ========== 18 WHY GREAT ==========
    s = new_light(prs)
    _kicker(s, "Эффект")
    _title(s, "Почему это автоматизация другого уровня")
    mets = [
        ("Часы → минуты", "Сбор нормативной базы проверки: не вкладки pravo.by, а утверждённый список и zip."),
        ("Кодекс → карточка", "Саммари даёт существенное для этой темы, а не дамп статей по порядку."),
        ("Память → цитата", "Ответ на норму опирается на файл кейса. Pipe сам дописывает источник."),
        ("Хаос → кейс", "Система записи: что утвердили, что скачали, какой URL, какой индекс."),
        ("Чат → артефакт", "Выход — Word и Excel, которые живут в рабочем процессе СВА, не в истории GPT."),
        ("Скрипт → агент", "Память, руки, мозг, стоп-кран, статусы. Не one-shot промпт."),
    ]
    for i, (h, d) in enumerate(mets):
        col, row = i % 3, i // 3
        x = Inches(0.4) + Inches(col * 4.25)
        y = Inches(1.4) + Inches(row * 2.55)
        c = _round(s, x, y, Inches(4.1), Inches(2.35), NAVY if row == 0 else WHITE, LINE if row else None, adj=0.05)
        title_c = GOLD if row == 0 else INK
        body_c = WHITE if row == 0 else SLATE
        _textbox(s, x + Inches(0.25), y + Inches(0.25), Inches(3.6), Inches(0.5), h, 16, title_c, bold=True)
        _textbox(s, x + Inches(0.25), y + Inches(0.85), Inches(3.6), Inches(1.25), d, 13, body_c)
    _footer(s, 18, TOTAL)

    # ========== 19 CLOSE ==========
    s = new_dark(prs)
    _kicker(s, "Итог", dark=True)
    _title(s, "Агент, которому можно доверить контур проверки", color=WHITE, size=26)
    left = _round(s, Inches(0.5), Inches(1.4), Inches(6.15), Inches(5.25), CARD, adj=0.04)
    tf = _tb(left, "Что уже можно показать", 15, GOLD, bold=True)
    tf.margin_left = Inches(0.28)
    tf.margin_top = Inches(0.22)
    for line in [
        "Один URL, модель «Аудитор» в списке.",
        "Список НПА нельзя скачать без «утверждаю».",
        "саммари / программа / гипотезы — файлы, не простыня в чате.",
        "вопрос … — с цитатой; нет в библиотеке — отказ.",
        "Клиентский текст не уходит в интернет.",
        "Свой фронт аудитору не нужен.",
    ]:
        _add_p(tf, "▸  " + line, 14, WHITE, space_before=12)

    right = _round(s, Inches(6.85), Inches(1.4), Inches(5.95), Inches(5.25), CARD, adj=0.04)
    tf = _tb(right, "Честная граница", 15, GOLD, bold=True)
    tf.margin_left = Inches(0.28)
    tf.margin_top = Inches(0.22)
    _add_p(tf, "Закрыт контур «собрать норму → сформулировать гипотезы».", 14, WHITE, space_before=10)
    _add_p(tf, "Дальше по тому же Pipe: подтверждение гипотез → черновики мнения и заключения. Затем — факты клиента, тесты, WP.", 14, WHITE, space_before=10)
    _add_p(tf, "Автоматизация сильна тем, что знает, где остановиться.", 15, GOLD2, bold=True, italic=True, space_before=18)
    _add_p(tf, "Откройте :3000 и начните с одной фразы.", 14, RGBColor(0xC5, 0xD4, 0xE0), space_before=14)
    _footer(s, 19, TOTAL, dark=True)

    out = Path(__file__).resolve().parent / "Аудитор_ИИ-агент_презентация.pptx"
    prs.save(out)
    print(out)


if __name__ == "__main__":
    build()
