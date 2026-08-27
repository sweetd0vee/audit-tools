"""Собрать презентацию для СВА / демо. Run: python docs/_build_presentation.py"""
from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pathlib import Path

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "Аудитор_ИИ-агент_презентация.pptx"

NAVY = RGBColor(0x0B, 0x1F, 0x33)
NAVY2 = RGBColor(0x12, 0x2C, 0x45)
NAVY3 = RGBColor(0x17, 0x33, 0x4C)
INK = RGBColor(0x14, 0x22, 0x33)
MUTED = RGBColor(0x3D, 0x4F, 0x61)
MUTED2 = RGBColor(0x6B, 0x7C, 0x8D)
GOLD = RGBColor(0xC4, 0x96, 0x2C)
LIGHT = RGBColor(0xD5, 0xE0, 0xEA)
CREAM = RGBColor(0xFF, 0xFC, 0xF7)
CREAM2 = RGBColor(0xF6, 0xF1, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 12
FOOTER = "Аудитор  ·  коробочный ИИ-агент внутренней проверки  ·  конфиденциально"


def _fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _style_run(run, size: int, color: RGBColor, *, bold: bool = False) -> None:
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Calibri"


def _set_p(p, text: str, size: int, color: RGBColor, *, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    p.alignment = align
    p.text = text
    if p.runs:
        _style_run(p.runs[0], size, color, bold=bold)


def _tf(shape, *, anchor=MSO_ANCHOR.TOP, margin: float = 0.08):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.anchor = anchor
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    return tf


def add_rect(slide, x, y, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _fill(sh, color)
    return sh


def add_round(slide, x, y, w, h, color, adj: float = 0.08):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _fill(sh, color)
    try:
        sh.adjustments[0] = adj
    except Exception:
        pass
    return sh


def add_tb(slide, x, y, w, h, text: str, size: int, color: RGBColor, *, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, margin=0.08):
    sh = slide.shapes.add_textbox(x, y, w, h)
    tf = _tf(sh, anchor=anchor, margin=margin)
    _set_p(tf.paragraphs[0], text, size, color, bold=bold, align=align)
    return sh


def add_lines(slide, x, y, w, h, lines: list[tuple], *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, margin=0.1, space_after=4):
    sh = slide.shapes.add_textbox(x, y, w, h)
    tf = _tf(sh, anchor=anchor, margin=margin)
    first = True
    for item in lines:
        text, size, color, bold = item[0], item[1], item[2], item[3] if len(item) > 3 else False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _set_p(p, text, size, color, bold=bold, align=align)
        p.space_after = Pt(space_after)
    return sh


def chrome_light(slide, page: int) -> None:
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CREAM)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), GOLD)
    add_rect(slide, 0, Inches(7.08), SLIDE_W, Inches(0.42), NAVY2)
    add_rect(slide, 0, Inches(7.08), SLIDE_W, Inches(0.02), GOLD)
    add_tb(slide, Inches(0.50), Inches(7.16), Inches(9.4), Inches(0.28), FOOTER, 10, RGBColor(0x9A, 0xB0, 0xC0), margin=0.02)
    add_tb(
        slide,
        Inches(11.20),
        Inches(7.16),
        Inches(1.70),
        Inches(0.28),
        f"{page}  /  {TOTAL}",
        10,
        RGBColor(0x9A, 0xB0, 0xC0),
        align=PP_ALIGN.RIGHT,
        margin=0.02,
    )


def chrome_dark(slide, page: int) -> None:
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, GOLD)
    add_rect(slide, 0, Inches(7.08), SLIDE_W, Inches(0.42), NAVY2)
    add_rect(slide, 0, Inches(7.08), SLIDE_W, Inches(0.02), GOLD)
    add_tb(slide, Inches(0.50), Inches(7.16), Inches(9.4), Inches(0.28), FOOTER, 10, RGBColor(0x9A, 0xB0, 0xC0), margin=0.02)
    add_tb(
        slide,
        Inches(11.20),
        Inches(7.16),
        Inches(1.70),
        Inches(0.28),
        f"{page}  /  {TOTAL}",
        10,
        RGBColor(0x9A, 0xB0, 0xC0),
        align=PP_ALIGN.RIGHT,
        margin=0.02,
    )


def header(slide, kicker: str, title: str, subtitle: str | None = None) -> None:
    add_rect(slide, Inches(0.55), Inches(0.35), Inches(0.18), Inches(0.18), GOLD)
    add_tb(slide, Inches(0.83), Inches(0.28), Inches(11.5), Inches(0.32), kicker, 11, GOLD, bold=True, margin=0.02)
    add_tb(slide, Inches(0.55), Inches(0.52), Inches(12.20), Inches(0.55), title, 26, NAVY, bold=True, margin=0.02)
    if subtitle:
        add_tb(slide, Inches(0.55), Inches(1.08), Inches(12.20), Inches(0.42), subtitle, 14, MUTED, margin=0.02)


def card_text(slide, x, y, w, h, title: str, body: str, *, accent: bool = False) -> None:
    bg = CREAM2 if accent else WHITE
    add_round(slide, x, y, w, h, bg)
    add_rect(slide, x, y, Inches(0.08), h, GOLD)
    add_lines(
        slide,
        x + Inches(0.16),
        y + Inches(0.10),
        w - Inches(0.28),
        h - Inches(0.18),
        [
            (title, 15, NAVY, True),
            (body, 13, MUTED, False),
        ],
        space_after=6,
    )


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def slide_title(prs: Presentation) -> None:
    s = blank(prs)
    chrome_dark(s, 1)
    add_tb(
        s,
        Inches(0.70),
        Inches(0.95),
        Inches(12.0),
        Inches(0.35),
        "ВНУТРЕННИЙ АУДИТ  ·  БАНК РБ  ·  НА МАШИНЕ БАНКА  ·  v0.0.1",
        13,
        LIGHT,
        bold=True,
        margin=0.02,
    )
    add_tb(s, Inches(0.70), Inches(1.55), Inches(12.0), Inches(0.95), "Аудитор", 54, WHITE, bold=True, margin=0.02)
    add_tb(
        s,
        Inches(0.70),
        Inches(2.50),
        Inches(11.8),
        Inches(1.05),
        "ИИ-агент внутренней проверки.\nПоставили — открыли чат — ведёте проверку.",
        22,
        LIGHT,
        margin=0.02,
    )
    cards = [
        ("Одна команда", "compose up → чат на :3000"),
        ("Готовый стек", "чат, модели, поиск — уже вместе"),
        ("Вы решаете", "модель предлагает, аудитор утверждает"),
        ("Не в облако", "данные остаются в контуре банка"),
    ]
    x0 = 0.70
    for i, (t, b) in enumerate(cards):
        x = Inches(x0 + i * 3.05)
        add_round(s, x, Inches(4.55), Inches(2.90), Inches(1.25), NAVY3, adj=0.10)
        add_lines(
            s,
            x,
            Inches(4.62),
            Inches(2.90),
            Inches(1.12),
            [(t, 16, GOLD, True), (b, 13, LIGHT, False)],
            space_after=6,
        )


def slide_toc(prs: Presentation) -> None:
    s = blank(prs)
    chrome_light(s, 2)
    header(s, "СОДЕРЖАНИЕ", "О чём эта презентация")
    items = [
        ("01", "Зачем", "Что болит на каждой проверке и чего модель не заменяет."),
        ("02", "Как работает", "От одной фразы в чате до черновика заключения в Word."),
        ("03", "Как устроено", "Три локальные модели, цитата из акта, два стоп-крана."),
        ("04", "Контур банка", "На машине банка, только официальные сайты РБ, без облака."),
        ("05", "Итог", "Что можно показать сегодня — и где человек обязателен."),
    ]
    y = 1.28
    for num, title, body in items:
        add_round(s, Inches(0.55), Inches(y), Inches(12.20), Inches(0.92), WHITE)
        add_tb(s, Inches(0.75), Inches(y + 0.18), Inches(0.90), Inches(0.56), num, 22, GOLD, bold=True, margin=0.02)
        add_lines(
            s,
            Inches(1.80),
            Inches(y + 0.14),
            Inches(10.50),
            Inches(0.66),
            [(title, 18, NAVY, True), (body, 14, MUTED, False)],
            space_after=2,
        )
        y += 1.02


def slide_why(prs: Presentation) -> None:
    s = blank(prs)
    chrome_light(s, 3)
    header(s, "ЗАЧЕМ", "Каждую проверку аудитор собирает заново")
    pains = [
        ("01   Вспомнить законы", "ГК, НК, инструкции НБРБ, акты Минфина и МНС — по теме этой недели."),
        ("02   Найти редакцию", "Актуальный текст на pravo.by / nbrb.by, не методичка с форума."),
        ("03   Вычитать кодекс", "Десятки страниц, чтобы вытащить 5–7 норм в рабочую бумагу."),
        ("04   Собрать бумаги СВА", "Программа, гипотезы, мнение, заключение — каждый раз с нуля."),
    ]
    y = 1.28
    for title, body in pains:
        add_round(s, Inches(0.45), Inches(y), Inches(6.20), Inches(1.22), WHITE)
        add_rect(s, Inches(0.45), Inches(y), Inches(0.10), Inches(1.22), GOLD)
        add_lines(
            s,
            Inches(0.70),
            Inches(y + 0.14),
            Inches(5.70),
            Inches(0.98),
            [(title, 15, NAVY, True), (body, 13, MUTED, False)],
            space_after=4,
        )
        y += 1.32
    add_round(s, Inches(6.90), Inches(1.28), Inches(5.95), Inches(5.30), NAVY)
    add_tb(s, Inches(7.15), Inches(1.50), Inches(5.50), Inches(0.35), "РАЗРЫВ", 12, GOLD, bold=True, margin=0.04)
    add_lines(
        s,
        Inches(7.15),
        Inches(1.95),
        Inches(5.50),
        Inches(4.35),
        [
            ("Чат без документов выдумывает статьи.", 16, WHITE, True),
            ("Документы без чата читаются слишком долго.", 16, WHITE, True),
            ("", 10, WHITE, False),
            ("Нужны оба: модель помогает найти, закон доказывает, аудитор подтверждает.", 15, LIGHT, False),
        ],
        space_after=10,
    )


def slide_product(prs: Presentation) -> None:
    s = blank(prs)
    chrome_light(s, 4)
    header(
        s,
        "ПРОДУКТ",
        "Коробка, а не «ещё один чат»",
        "Берём готовый чат и модели. Своё пишем только там, где рынку нечего дать аудиту банка РБ.",
    )
    cols = [
        (
            "Берём готовое",
            "Чат с входом и файлами.\nЛокальные модели.\nПоиск по сайтам РБ.\nWord и Excel по шаблону.\nЗапуск одной командой.",
        ),
        (
            "Пишем своё",
            "Проверка как дело на диске.\nСкачивание законов РБ.\nНарезка текста по статьям.\nПромпты методолога ВА.\nСборка черновиков документов.",
        ),
        (
            "Не пишем вовсе",
            "Свой чат и свой сайт.\nСвою нейросеть.\nСвой поисковик.\nРедактор рабочих бумаг.\nПодпись «за аудитора».",
        ),
    ]
    x0 = 0.45
    for i, (t, b) in enumerate(cols):
        x = Inches(x0 + i * 4.25)
        add_round(s, x, Inches(1.62), Inches(4.10), Inches(3.35), WHITE)
        add_tb(s, x + Inches(0.18), Inches(1.76), Inches(3.74), Inches(0.40), t, 16, GOLD, bold=True)
        add_tb(s, x + Inches(0.18), Inches(2.22), Inches(3.74), Inches(2.50), b, 14, INK)
    add_round(s, Inches(0.45), Inches(5.12), Inches(12.40), Inches(1.70), CREAM2)
    add_lines(
        s,
        Inches(0.65),
        Inches(5.22),
        Inches(12.00),
        Inches(1.50),
        [
            (
                "«Модель ускоряет поиск нормы и черновик. Решение и подпись — за аудитором. Цитату сверяют с файлом в деле.»",
                15,
                NAVY,
                True,
            ),
            ("Ускоряет сбор законов и бумаг. Не ставит вывод «нарушение». Не подписывает заключение.", 13, MUTED, False),
        ],
        space_after=6,
    )


def slide_session(prs: Presentation) -> None:
    s = blank(prs)
    chrome_light(s, 5)
    header(
        s,
        "КАК РАБОТАЕТ",
        "Один чат — одна проверка до черновика",
        "Тема: аренда коммерческой недвижимости — валюта, НДС, договоры с нерезидентами.",
    )
    add_round(s, Inches(0.45), Inches(1.58), Inches(6.05), Inches(2.55), WHITE)
    add_tb(s, Inches(0.65), Inches(1.70), Inches(5.65), Inches(0.35), "РАНЬШЕ", 13, GOLD, bold=True)
    add_tb(
        s,
        Inches(0.65),
        Inches(2.08),
        Inches(5.65),
        Inches(1.90),
        "Вспоминаю акты по памяти и старым бумагам.\nИщу редакции вкладками на pravo.by.\nЧитаю ГК и НК «на всякий случай».\nПрограмму, мнение и заключение пишу с нуля.",
        14,
        INK,
    )
    add_round(s, Inches(6.80), Inches(1.58), Inches(6.05), Inches(2.55), CREAM2)
    add_tb(s, Inches(7.00), Inches(1.70), Inches(5.65), Inches(0.35), "ТЕПЕРЬ", 13, GOLD, bold=True)
    add_tb(
        s,
        Inches(7.00),
        Inches(2.08),
        Inches(5.65),
        Inches(1.90),
        "Одной фразой описываю проверку в чате.\nУтверждаю список законов — без этого ничего не качается.\nЧитаю Word-обзор вместо кодекса целиком.\nГипотезы в Excel → черновик мнения и заключения.",
        14,
        INK,
    )
    steps = [
        ("1", "Описать", "Проверка аренды…"),
        ("2", "Утвердить", "утверждаю 1, 2, 4"),
        ("3", "Читать", "саммари / программа"),
        ("4", "Планировать", "гипотезы → Excel"),
        ("5", "Подтвердить", "утверждаю гипотезы"),
        ("6", "Мнение", "раздел I, 2–4 стр."),
        ("7", "Заключение", "полный черновик"),
        ("8", "Спросить", "вопрос Какой срок…"),
    ]
    x0 = 0.45
    for i, (n, t, b) in enumerate(steps):
        x = Inches(x0 + i * 1.58)
        add_round(s, x, Inches(4.30), Inches(1.50), Inches(2.50), WHITE, adj=0.12)
        add_tb(s, x, Inches(4.42), Inches(1.50), Inches(0.40), n, 18, GOLD, bold=True, align=PP_ALIGN.CENTER)
        add_tb(s, x, Inches(4.85), Inches(1.50), Inches(0.70), t, 13, NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_tb(s, x + Inches(0.06), Inches(5.55), Inches(1.38), Inches(1.05), b, 11, MUTED, align=PP_ALIGN.CENTER)


def slide_outputs(prs: Presentation) -> None:
    s = blank(prs)
    chrome_light(s, 6)
    header(
        s,
        "ЧТО ПОЛУЧАЕТ АУДИТОР",
        "Файлы, которые уносите из чата",
        "В мнение и заключение попадают только гипотезы, которые вы подтвердили.",
    )
    items = [
        (".zip", "Библиотека законов", "Официальные тексты, которые вы утвердили — не то, что модель «вспомнила»."),
        (".docx", "Обзор актов", "Зачем каждый закон и какие статьи смотреть. Со ссылками на фрагменты."),
        (".docx", "Программа", "Черновик программы проверки: цель, риски, что запросить, как проверить."),
        (".xlsx", "Гипотезы", "Чеклист 8–10 гипотез: что смотреть и на что опираться."),
        (".docx", "Мнение", "Раздел I для руководства банка. Черновик для правки, не подпись."),
        (".docx", "Заключение", "Титул, содержание, мнение, наблюдения. Не акт СВА."),
    ]
    for i, (ext, title, body) in enumerate(items):
        col, row = i % 3, i // 3
        x = Inches(0.45 + col * 4.25)
        y = Inches(1.58 + row * 2.55)
        add_round(s, x, y, Inches(4.10), Inches(2.40), WHITE)
        add_tb(s, x + Inches(0.20), y + Inches(0.18), Inches(3.70), Inches(0.35), ext, 14, GOLD, bold=True)
        add_tb(s, x + Inches(0.20), y + Inches(0.52), Inches(3.70), Inches(0.45), title, 16, NAVY, bold=True)
        add_tb(s, x + Inches(0.20), y + Inches(1.05), Inches(3.70), Inches(1.10), body, 13, MUTED)


def slide_arch(prs: Presentation) -> None:
    s = blank(prs)
    chrome_light(s, 7)
    header(s, "КАК УСТРОЕНО", "Сборка готового, а не свой чат с нуля")
    layers = [
        ("Чат", "Одно окно для аудитора. Вход, файлы, статусы. Свой сайт не пишем."),
        ("Агент «Аудитор»", "Понимает обычные фразы: «утверждаю», «саммари», «вопрос …». Дальше вызывает сервер."),
        ("Сервер проверки", "Дело на диске: какие законы утверждены, что скачано, какие черновики собраны."),
        ("На машине банка", "Модели локально. Поиск только по официальным сайтам РБ. Файлы дела — на диске."),
    ]
    y = 1.30
    for title, body in layers:
        add_round(s, Inches(0.50), Inches(y), Inches(12.30), Inches(1.25), WHITE)
        add_rect(s, Inches(0.50), Inches(y), Inches(0.10), Inches(1.25), GOLD)
        add_lines(
            s,
            Inches(0.80),
            Inches(y + 0.22),
            Inches(11.70),
            Inches(0.90),
            [(title, 18, NAVY, True), (body, 15, MUTED, False)],
            space_after=4,
        )
        y += 1.38


def slide_cite(prs: Presentation) -> None:
    s = blank(prs)
    chrome_light(s, 8)
    header(
        s,
        "ЦИТАТА",
        "Модель думает. Документы доказывают.",
        "На «вопрос …» в ответ кладутся куски утверждённых актов. Нет куска — нет статьи.",
    )
    roles = [
        ("Чат", "Пишет список законов, обзор, программу, мнение и заключение."),
        ("Поиск по смыслу", "Находит нужную статью в библиотеке этой проверки."),
        ("Перепроверка", "Отсекает похожие, но не те фрагменты — чтобы не перепутать статьи."),
    ]
    for i, (t, b) in enumerate(roles):
        x = Inches(0.45 + i * 4.25)
        add_round(s, x, Inches(1.58), Inches(4.10), Inches(1.70), WHITE)
        add_tb(s, x + Inches(0.18), Inches(1.70), Inches(3.74), Inches(0.40), t, 15, GOLD, bold=True)
        add_tb(s, x + Inches(0.18), Inches(2.12), Inches(3.74), Inches(0.95), b, 13, INK)
    add_round(s, Inches(0.45), Inches(3.45), Inches(7.55), Inches(3.35), WHITE)
    add_tb(s, Inches(0.65), Inches(3.55), Inches(7.15), Inches(0.35), "Как отвечает на вопрос", 15, NAVY, bold=True)
    add_tb(
        s,
        Inches(0.65),
        Inches(3.95),
        Inches(7.15),
        Inches(2.65),
        "1.  Пишете «вопрос …». Без этого слова — просто разговор, не норма.\n"
        "2.  Система ищет куски в законах, которые вы утвердили.\n"
        "3.  Модель видит только эти куски — не «память» и не интернет.\n"
        "4.  Куска нет — отказ. Статью из головы не подставляет.\n"
        "5.  В ответе блок «Откуда» — сверяете с файлом в деле.",
        14,
        INK,
    )
    add_round(s, Inches(8.20), Inches(3.45), Inches(4.65), Inches(3.35), CREAM2)
    add_tb(s, Inches(8.40), Inches(3.55), Inches(4.25), Inches(0.35), "Закон ≠ цифра клиента", 15, NAVY, bold=True)
    add_tb(
        s,
        Inches(8.40),
        Inches(4.00),
        Inches(4.25),
        Inches(2.55),
        "Закон — из библиотеки проверки, со ссылкой на статью.\n\n"
        "Цифры клиента (договоры, выгрузки) — следующий этап. Excel в базу законов не кладём.",
        14,
        INK,
    )


def slide_hitl(prs: Presentation) -> None:
    s = blank(prs)
    chrome_light(s, 9)
    header(s, "АУДИТОР РЕШАЕТ", "Два стоп-крана: законы и гипотезы")
    cranes = [
        (
            "1. Законы",
            "Модель предлагает список актов.\nПока не напишете «утверждаю 1, 2, 4» — ничего не скачивается.",
        ),
        (
            "2. Гипотезы",
            "Excel готов. Пока не напишете «утверждаю гипотезы 1, 3, 5» — нет мнения и заключения.",
        ),
    ]
    for i, (t, b) in enumerate(cranes):
        x = Inches(0.45 + i * 6.40)
        add_round(s, x, Inches(1.30), Inches(6.20), Inches(2.15), WHITE)
        add_rect(s, x, Inches(1.30), Inches(0.10), Inches(2.15), GOLD)
        add_lines(
            s,
            x + Inches(0.28),
            Inches(1.48),
            Inches(5.70),
            Inches(1.85),
            [(t, 18, NAVY, True), (b, 15, MUTED, False)],
            space_after=8,
        )
    add_round(s, Inches(0.45), Inches(3.65), Inches(6.20), Inches(3.15), WHITE)
    add_tb(s, Inches(0.65), Inches(3.80), Inches(5.80), Inches(0.40), "Где модель молчит", 16, GOLD, bold=True)
    add_tb(
        s,
        Inches(0.65),
        Inches(4.28),
        Inches(5.80),
        Inches(2.30),
        "утверждаю, скачай, документы, статус, помощь — это код, не нейросеть.\n\nВыдумать статью на этом шаге нельзя: качается только то, что вы выбрали.",
        14,
        INK,
    )
    add_round(s, Inches(6.85), Inches(3.65), Inches(6.00), Inches(3.15), CREAM2)
    add_tb(s, Inches(7.05), Inches(3.80), Inches(5.60), Inches(0.40), "Где модель говорит", 16, GOLD, bold=True)
    add_tb(
        s,
        Inches(7.05),
        Inches(4.28),
        Inches(5.60),
        Inches(2.30),
        "Список актов, тексты Word и Excel, ответ на вопрос по найденным кускам, обычный разговор в чате.\n\nЧерновик можно править. Подпись — ваша.",
        14,
        INK,
    )


def slide_privacy(prs: Presentation) -> None:
    s = blank(prs)
    chrome_light(s, 10)
    header(s, "КОНТУР БАНКА", "Данные не уходят в облако — так задумано")
    items = [
        ("Модели на месте", "Чат и поиск крутятся на машине банка. Формулировки проверки не уходят в интернет."),
        ("Только сайты РБ", "Качаем pravo.by, nbrb.by и другие официальные домены. Внешний поиск выключен."),
        ("Клиент не в поиск", "Выгрузки, договоры и суммы в поисковик не передаются. Поиск в чате выключен."),
        ("Две полки", "Закон и цифры клиента не смешиваются. Excel клиента в базу законов класть нельзя."),
        ("Только эта машина", "Чат и сервер слушают localhost. С соседнего компьютера контур не открывается."),
        ("След вопроса", "Какой вопрос задали и из какого фрагмента ответили — пишется в журнал дела."),
    ]
    for i, (t, b) in enumerate(items):
        col, row = i % 3, i // 3
        x = Inches(0.45 + col * 4.25)
        y = Inches(1.30 + row * 2.70)
        add_round(s, x, y, Inches(4.10), Inches(2.50), WHITE)
        add_tb(s, x + Inches(0.20), y + Inches(0.22), Inches(3.70), Inches(0.50), t, 16, NAVY, bold=True)
        add_tb(s, x + Inches(0.20), y + Inches(0.80), Inches(3.70), Inches(1.45), b, 14, MUTED)


def slide_box(prs: Presentation) -> None:
    s = blank(prs)
    chrome_light(s, 11)
    header(
        s,
        "ПОСТАВКА",
        "Коробка, а не стенд разработчика",
        "Аудитору не нужно выбирать модели и включать поиск. Открыл чат — работает.",
    )
    add_round(s, Inches(0.45), Inches(1.55), Inches(12.40), Inches(0.85), NAVY)
    add_tb(
        s,
        Inches(0.65),
        Inches(1.70),
        Inches(12.00),
        Inches(0.55),
        "docker compose up -d --build     →     http://localhost:3000     →     модель «Аудитор»",
        16,
        WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    cols = [
        (
            "Уже в коробке",
            "Чат, сервер и поиск — одна команда.\nМодели локально, на хосте банка.\nПоиск в чате выключен.\nАгент ставится при запуске (нужен ключ администратора).",
        ),
        (
            "Качество сразу",
            "Промпты методолога ВА банка РБ.\nНет фрагмента — отказ, не статья «из памяти».\nНормы РФ / ЕС / IFRS сами не подставляет.\nДва стоп-крана: законы и гипотезы.",
        ),
        (
            "Что это даёт",
            "Один адрес для аудитора.\nНе облачный сервис с файлами банка снаружи.\nМожно поставить в контур СВА.\nСледующая проверка — тот же чат, новое окно.",
        ),
    ]
    for i, (t, b) in enumerate(cols):
        x = Inches(0.45 + i * 4.25)
        add_round(s, x, Inches(2.58), Inches(4.10), Inches(4.20), WHITE)
        add_tb(s, x + Inches(0.20), Inches(2.75), Inches(3.70), Inches(0.45), t, 16, GOLD, bold=True)
        add_tb(s, x + Inches(0.20), Inches(3.25), Inches(3.70), Inches(3.25), b, 14, INK)


def slide_close(prs: Presentation) -> None:
    s = blank(prs)
    chrome_light(s, 12)
    header(s, "ИТОГ", "Агент, которому можно доверить контур проверки")
    add_round(s, Inches(0.45), Inches(1.30), Inches(6.15), Inches(4.15), WHITE)
    add_tb(s, Inches(0.65), Inches(1.45), Inches(5.75), Inches(0.40), "Что уже можно показать", 16, GOLD, bold=True)
    add_tb(
        s,
        Inches(0.65),
        Inches(1.95),
        Inches(5.75),
        Inches(3.25),
        "Один адрес, модель «Аудитор» в списке.\n\n"
        "Без «утверждаю» законы не качаются и не идут в заключение.\n\n"
        "Обзор, программа, гипотезы, мнение, заключение — файлы.\n\n"
        "«вопрос …» — с цитатой. Нет в библиотеке — отказ.\n\n"
        "Клиентский текст не уходит в интернет.",
        14,
        INK,
    )
    add_round(s, Inches(6.80), Inches(1.30), Inches(6.05), Inches(4.15), CREAM2)
    add_tb(s, Inches(7.00), Inches(1.45), Inches(5.65), Inches(0.40), "Честная граница", 16, GOLD, bold=True)
    add_tb(
        s,
        Inches(7.00),
        Inches(1.95),
        Inches(5.65),
        Inches(3.25),
        "Закрыт контур «норма → гипотезы → черновик мнения → черновик заключения».\n\n"
        "Это черновики для правки. Не подписанный акт СВА и не суждение «за аудитора».\n\n"
        "Дальше по тому же агенту: факты клиента, тесты, рабочие бумаги.",
        14,
        INK,
    )
    add_round(s, Inches(0.45), Inches(5.60), Inches(12.40), Inches(1.20), NAVY)
    add_tb(
        s,
        Inches(0.65),
        Inches(5.72),
        Inches(12.00),
        Inches(0.95),
        "Автоматизация сильна тем, что знает, где остановиться.\nОткройте :3000 и начните с одной фразы.",
        16,
        WHITE,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def set_props(prs: Presentation) -> None:
    prs.core_properties.title = "Аудитор — коробочный ИИ-агент внутренней проверки"
    prs.core_properties.subject = "v0.0.1 · внутренний аудит банка РБ · on-prem"
    prs.core_properties.author = "Audit Tools"
    prs.core_properties.comments = "generated using python-pptx"


def main() -> None:
    prs = new_prs()
    slide_title(prs)
    slide_toc(prs)
    slide_why(prs)
    slide_product(prs)
    slide_session(prs)
    slide_outputs(prs)
    slide_arch(prs)
    slide_cite(prs)
    slide_hitl(prs)
    slide_privacy(prs)
    slide_box(prs)
    slide_close(prs)
    set_props(prs)
    prs.save(OUT)
    print(f"Wrote {OUT} ({TOTAL} slides)")


if __name__ == "__main__":
    main()
